import os
import docx
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def extract_text_from_docx(filepath: str) -> str:
    """Extracts all paragraph text from a DOCX file."""
    try:
        doc = docx.Document(filepath)
        text_blocks = []
        for para in doc.paragraphs:
            if para.text.strip():
                text_blocks.append(para.text.strip())
        return "\n\n".join(text_blocks)
    except Exception as e:
        print(f"Error parsing DOCX {filepath}: {e}")
        return ""

def extract_text_from_pptx(filepath: str) -> str:
    """Extracts all text shapes from a PPTX file."""
    try:
        prs = Presentation(filepath)
        text_blocks = []
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            if slide_text:
                text_blocks.append(f"Slide {i+1}:\n" + "\n".join(slide_text))
        return "\n\n---\n\n".join(text_blocks)
    except Exception as e:
        print(f"Error parsing PPTX {filepath}: {e}")
        return ""

def extract_images_from_docx(filepath: str, output_dir: str) -> dict:
    """Extracts all images from a DOCX file and saves them to output_dir."""
    image_map = {}
    try:
        doc = docx.Document(filepath)
        os.makedirs(output_dir, exist_ok=True)
        
        images = []
        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                img_ext = rel.target_ref.split('.')[-1]
                img_bytes = rel.target_part.blob
                
                img_name = f"image_{len(images) + 1}.{img_ext}"
                img_path = os.path.join(output_dir, img_name)
                
                with open(img_path, "wb") as f:
                    f.write(img_bytes)
                    
                images.append(img_path)
        
        if images:
            # DOCX doesn't have strict pagination in raw XML, map all images to "1"
            image_map["1"] = images
            
    except Exception as e:
        print(f"Error extracting images from DOCX {filepath}: {e}")
        
    return image_map

def extract_images_from_pptx(filepath: str, output_dir: str) -> dict:
    """Extracts all image shapes from a PPTX file and saves them to output_dir."""
    image_map = {}
    try:
        prs = Presentation(filepath)
        os.makedirs(output_dir, exist_ok=True)
        
        img_count = 0
        for i, slide in enumerate(prs.slides):
            slide_images = []
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    img_bytes = shape.image.blob
                    img_ext = shape.image.ext
                    
                    img_count += 1
                    img_name = f"slide_{i+1}_image_{img_count}.{img_ext}"
                    img_path = os.path.join(output_dir, img_name)
                    
                    with open(img_path, "wb") as f:
                        f.write(img_bytes)
                        
                    slide_images.append(img_path)
            
            if slide_images:
                # Map images to the slide number (1-indexed)
                image_map[str(i+1)] = slide_images
                
    except Exception as e:
        print(f"Error extracting images from PPTX {filepath}: {e}")
        
    return image_map
